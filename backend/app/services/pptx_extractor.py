"""
PowerPoint text extraction service.

This module provides functions to extract text content from PowerPoint files (.pptx).
"""

import io
import logging
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches

logger = logging.getLogger(__name__)


def extract_text_from_pptx(content: bytes) -> tuple[str, List[Dict]]:
    """
    Extract text from a PowerPoint file.

    Args:
        content: PowerPoint file content as bytes

    Returns:
        Tuple of (full_text, slides_data)
        - full_text: All text concatenated with slide markers
        - slides_data: List of dicts with slide number, text, and notes

    Raises:
        ValueError: If the content is not a valid PowerPoint file
    """
    try:
        # Load presentation from bytes
        prs = Presentation(io.BytesIO(content))
    except Exception as e:
        logger.error(f"Failed to parse PowerPoint file: {e}")
        raise ValueError(f"PowerPointファイルの読み込みに失敗しました: {str(e)}")

    slides_data: List[Dict] = []
    full_text_parts: List[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: List[str] = []
        notes_text = ""

        # Extract text from all shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text = shape.text.strip()
                if text:
                    slide_texts.append(text)

            # Handle tables
            if shape.has_table:
                table_texts = _extract_table_text(shape.table)
                slide_texts.extend(table_texts)

            # Handle grouped shapes
            if hasattr(shape, "shapes"):
                group_texts = _extract_group_text(shape)
                slide_texts.extend(group_texts)

        # Extract notes
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame and notes_frame.text:
                notes_text = notes_frame.text.strip()

        slide_text = "\n".join(slide_texts)

        slides_data.append(
            {
                "slide": slide_num,
                "text": slide_text,
                "notes": notes_text,
            }
        )

        # Build full text with slide markers
        full_text_parts.append(f"[スライド {slide_num}]")
        if slide_text:
            full_text_parts.append(slide_text)
        if notes_text:
            full_text_parts.append(f"（ノート: {notes_text}）")
        full_text_parts.append("")  # Empty line between slides

    full_text = "\n".join(full_text_parts)
    page_count = len(prs.slides)

    logger.info(
        f"Extracted text from {page_count} slides, total {len(full_text)} chars"
    )

    return full_text, slides_data


def _extract_table_text(table) -> List[str]:
    """Extract text from a table shape."""
    texts = []
    for row in table.rows:
        row_texts = []
        for cell in row.cells:
            if cell.text:
                row_texts.append(cell.text.strip())
        if row_texts:
            texts.append(" | ".join(row_texts))
    return texts


def _extract_group_text(group_shape) -> List[str]:
    """Extract text from grouped shapes recursively."""
    texts = []
    for shape in group_shape.shapes:
        if hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
            if text:
                texts.append(text)
        if hasattr(shape, "shapes"):
            texts.extend(_extract_group_text(shape))
    return texts


def get_slide_count(content: bytes) -> int:
    """
    Get the number of slides in a PowerPoint file.

    Args:
        content: PowerPoint file content as bytes

    Returns:
        Number of slides
    """
    try:
        prs = Presentation(io.BytesIO(content))
        return len(prs.slides)
    except Exception as e:
        logger.error(f"Failed to get slide count: {e}")
        return 0
