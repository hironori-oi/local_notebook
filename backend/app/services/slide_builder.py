"""
Slide Builder service for creating PowerPoint files.

This module provides functions to build .pptx files from slide data.
"""
import io
import logging
from typing import List, Dict, Optional, Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

# Default style settings
DEFAULT_STYLE = {
    "colors": {
        "primary": "#1a73e8",
        "secondary": "#5f6368",
        "accent": "#ea4335",
        "background": "#ffffff",
        "text": "#202124",
        "light_text": "#5f6368",
    },
    "fonts": {
        "title": "Yu Gothic UI",
        "body": "Yu Gothic UI",
    },
    "sizes": {
        "title": 44,
        "subtitle": 28,
        "section_title": 40,
        "slide_title": 32,
        "body": 20,
        "bullet": 18,
        "notes": 12,
    },
}


def hex_to_rgb(hex_color: str) -> RGBColor:
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


class SlideBuilder:
    """
    Build PowerPoint presentations from slide data.

    Supports:
    - Custom templates
    - Custom style settings
    - Various slide types (title, content, section, conclusion)
    """

    def __init__(
        self,
        template_path: Optional[str] = None,
        style: Optional[Dict[str, Any]] = None,
    ):
        """
        Initialize SlideBuilder.

        Args:
            template_path: Path to .pptx template file
            style: Custom style settings dict
        """
        if template_path:
            self.prs = Presentation(template_path)
            self._use_template = True
        else:
            self.prs = Presentation()
            # Set slide size to 16:9
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
            self._use_template = False

        # Merge custom style with defaults
        self.style = DEFAULT_STYLE.copy()
        if style:
            for key in style:
                if isinstance(style[key], dict) and key in self.style:
                    self.style[key] = {**self.style[key], **style[key]}
                else:
                    self.style[key] = style[key]

    def add_title_slide(
        self,
        title: str,
        subtitle: Optional[str] = None,
        speaker_notes: Optional[str] = None,
    ) -> None:
        """
        Add a title slide.

        Args:
            title: Main title
            subtitle: Optional subtitle
            speaker_notes: Optional speaker notes
        """
        if self._use_template and len(self.prs.slide_layouts) > 0:
            # Try to use template's title layout
            layout = self.prs.slide_layouts[0]
        else:
            layout = self.prs.slide_layouts[6]  # Blank layout

        slide = self.prs.slides.add_slide(layout)

        # Add background color
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = hex_to_rgb(self.style["colors"]["primary"])
        background.line.fill.background()

        # Move background to back
        spTree = slide.shapes._spTree
        sp = background._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(self.style["sizes"]["title"])
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.font.name = self.style["fonts"]["title"]
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle if provided
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2),
                Inches(12.333), Inches(1)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = subtitle
            subtitle_para.font.size = Pt(self.style["sizes"]["subtitle"])
            subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
            subtitle_para.font.name = self.style["fonts"]["body"]
            subtitle_para.alignment = PP_ALIGN.CENTER

        # Add speaker notes
        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def add_section_slide(
        self,
        title: str,
        speaker_notes: Optional[str] = None,
    ) -> None:
        """
        Add a section divider slide.

        Args:
            title: Section title
            speaker_notes: Optional speaker notes
        """
        layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(layout)

        # Add accent bar
        accent_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(3.2),
            Inches(13.333), Inches(0.1)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = hex_to_rgb(self.style["colors"]["primary"])
        accent_bar.line.fill.background()

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5),
            Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(self.style["sizes"]["section_title"])
        title_para.font.bold = True
        title_para.font.color.rgb = hex_to_rgb(self.style["colors"]["text"])
        title_para.font.name = self.style["fonts"]["title"]
        title_para.alignment = PP_ALIGN.CENTER

        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def add_content_slide(
        self,
        title: str,
        bullets: Optional[List[str]] = None,
        details: Optional[str] = None,
        speaker_notes: Optional[str] = None,
    ) -> None:
        """
        Add a content slide with bullets.

        Args:
            title: Slide title
            bullets: List of bullet points
            details: Optional detailed text
            speaker_notes: Optional speaker notes
        """
        layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(layout)

        # Add title bar background
        title_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(13.333), Inches(1.2)
        )
        title_bg.fill.solid()
        title_bg.fill.fore_color.rgb = hex_to_rgb(self.style["colors"]["primary"])
        title_bg.line.fill.background()

        # Move to back
        spTree = slide.shapes._spTree
        sp = title_bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        # Add title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3),
            Inches(12.333), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(self.style["sizes"]["slide_title"])
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.font.name = self.style["fonts"]["title"]

        # Add bullets
        if bullets:
            content_box = slide.shapes.add_textbox(
                Inches(0.75), Inches(1.5),
                Inches(11.833), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, bullet in enumerate(bullets):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()

                para.text = f"• {bullet}"
                para.font.size = Pt(self.style["sizes"]["bullet"])
                para.font.color.rgb = hex_to_rgb(self.style["colors"]["text"])
                para.font.name = self.style["fonts"]["body"]
                para.space_after = Pt(12)

        # Add details if provided
        if details:
            details_top = Inches(1.5 + (len(bullets or []) * 0.5))
            if details_top > Inches(5):
                details_top = Inches(5)

            details_box = slide.shapes.add_textbox(
                Inches(0.75), details_top,
                Inches(11.833), Inches(2)
            )
            details_frame = details_box.text_frame
            details_frame.word_wrap = True
            details_para = details_frame.paragraphs[0]
            details_para.text = details
            details_para.font.size = Pt(self.style["sizes"]["body"])
            details_para.font.color.rgb = hex_to_rgb(self.style["colors"]["light_text"])
            details_para.font.name = self.style["fonts"]["body"]

        if speaker_notes:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = speaker_notes

    def add_conclusion_slide(
        self,
        title: str,
        bullets: Optional[List[str]] = None,
        speaker_notes: Optional[str] = None,
    ) -> None:
        """
        Add a conclusion slide.

        Args:
            title: Slide title
            bullets: Summary points
            speaker_notes: Optional speaker notes
        """
        # Use same layout as content slide but with different styling
        self.add_content_slide(title, bullets, None, speaker_notes)

    def build_from_slides(self, slides: List[Dict]) -> bytes:
        """
        Build presentation from slide data list.

        Args:
            slides: List of slide dicts with structure:
                {
                    "slide_type": "title|section|content|conclusion",
                    "title": "...",
                    "content": {"bullets": [...], "subtitle": "...", "details": "..."},
                    "speaker_notes": "..."
                }

        Returns:
            PowerPoint file as bytes
        """
        for slide_data in slides:
            slide_type = slide_data.get("slide_type", "content")
            title = slide_data.get("title", "")
            content = slide_data.get("content", {})
            speaker_notes = slide_data.get("speaker_notes", "")

            if slide_type == "title":
                self.add_title_slide(
                    title=title,
                    subtitle=content.get("subtitle"),
                    speaker_notes=speaker_notes,
                )
            elif slide_type == "section":
                self.add_section_slide(
                    title=title,
                    speaker_notes=speaker_notes,
                )
            elif slide_type == "conclusion":
                self.add_conclusion_slide(
                    title=title,
                    bullets=content.get("bullets", []),
                    speaker_notes=speaker_notes,
                )
            else:  # content
                self.add_content_slide(
                    title=title,
                    bullets=content.get("bullets", []),
                    details=content.get("details"),
                    speaker_notes=speaker_notes,
                )

        # Save to bytes
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output.getvalue()

    def save(self, output_path: str) -> None:
        """
        Save presentation to file.

        Args:
            output_path: Output file path
        """
        self.prs.save(output_path)


def build_pptx_from_project(
    slides: List[Dict],
    template_path: Optional[str] = None,
    style: Optional[Dict] = None,
) -> bytes:
    """
    Convenience function to build PPTX from slide data.

    Args:
        slides: List of slide dicts
        template_path: Optional template file path
        style: Optional custom style settings

    Returns:
        PowerPoint file as bytes
    """
    builder = SlideBuilder(template_path=template_path, style=style)
    return builder.build_from_slides(slides)
